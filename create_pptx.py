import datetime
import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE  # Add MSO_AUTO_SIZE here
from pptx.dml.color import RGBColor
import textwrap
import io
from PIL import Image


class PowerPointHelper:
    def __init__(self, show_pptx):
        self.show_pptx = show_pptx
        self.presentation = Presentation()

    def create_presentation(self):
        if not self.show_pptx["Verses"]:
            return io.BytesIO()  # Return an empty BytesIO object

        for verse in self.show_pptx["Verses"]:
            self.create_slide(verse)

        output = io.BytesIO()
        self.presentation.save(output)
        output.seek(0)
        return output

    def create_slide(self, verse):
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.set_background(slide)
       
        # 1. Bold
        # 2. Italic
        # 3. Bold and Italic
        # 4. Underline
        # 5. Underline and Bold
        # 6. Underline and Italic
        # 7. Bold and Italic and Underline
        if self.show_pptx["Config"]:
            font_family = self.show_pptx["Config"].get("FontFamily", "Arial")
            font_size = Pt(self.show_pptx["Config"].get("FontSize", 40))
            font_style = self.show_pptx["Config"].get("FontStyle", 1)
            if self.show_pptx["Config"].get("Color"):
                font_color = RGBColor(
                    self.show_pptx["Config"]["Color"]["R"],
                    self.show_pptx["Config"]["Color"]["G"],
                    self.show_pptx["Config"]["Color"]["B"],
                )
            else:
                font_color = RGBColor(255, 255, 255)

            text_align = getattr(
                PP_ALIGN, self.show_pptx["Config"].get("TextAlign", "CENTER").upper()
            )
        else:
            font_family = "Arial"
            font_size = Pt(40)
            font_style = 1
            font_color = RGBColor(255, 255, 255)
            text_align = PP_ALIGN.CENTER

        # add Content
        content = f"{verse['label']}.  {verse['content']}"
        textbox = slide.shapes.add_textbox(
            Inches(0.2), Inches(1), Inches(9), Inches(8)  # Adjust height to allow for expansion
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Enable auto-sizing

        p = text_frame.add_paragraph()
        p.text = content
        p.font.name = font_family
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = font_style in [1, 3, 5, 7]
        p.font.italic = font_style in [2, 3, 6, 7]
        p.font.underline = font_style in [4, 5, 6, 7]
        p.alignment = text_align

        #  add Title
        title = f"{self.show_pptx['BookName']} {self.show_pptx['ChapterNumber']}:{verse['label']}"
        textbox = slide.shapes.add_textbox(Inches(1), Inches(-0.3), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = title
        p.font.name = font_family
        p.font.size = font_size
        p.font.bold = font_style in [1, 3, 5, 7]
        p.font.italic = font_style in [2, 3, 6, 7]
        p.font.underline = font_style in [4, 5, 6, 7]
        p.font.color.rgb = RGBColor(150, 150, 57)
        p.alignment = PP_ALIGN.CENTER

    def set_background(self, slide):
        if self.show_pptx["Config"] and self.show_pptx["Config"].get("ImagePath") and os.path.isfile(self.show_pptx["Config"]["ImagePath"]) and self.show_pptx["Config"]["ImagePath"] !="Choose Image":
            slide.shapes.add_picture(
                self.show_pptx["Config"]["ImagePath"],
                0,
                0,
                self.presentation.slide_width,
                self.presentation.slide_height,
            )
            return

        slide_width =  1920
        slide_height = 1080
        img = Image.new("RGB", (int(slide_width), int(slide_height)), color="black")
        img_path = "temp_background.png"
        img.save(img_path)
        slide.shapes.add_picture(
            img_path,
            0,
            0,
            self.presentation.slide_width,
            self.presentation.slide_height,
        )
        os.remove(img_path)


if __name__ == "__main__":
    print("Creating PowerPoint presentation...")
    try:
        # Read input JSON data from stdin
        print("Reading input data...")
        input_data = sys.stdin.read()
        if not input_data:
            print("No input provided")
        else:
            # Proceed with the rest of your code
            pass
        print("Input data: ", input_data)
        print("Input data read successfully.")
        show_pptx = json.loads(input_data)
        print("Converted JSON data to Python object.")

        helper = PowerPointHelper(show_pptx)
        print("Created PowerPoint presentation.")

        output = helper.create_presentation()
        print("Presentation created successfully.")
        with open(show_pptx["FilePath"], "wb") as f:
            f.write(output.getbuffer())
        print(f"Presentation saved to {show_pptx['FilePath']}")
        os.startfile(show_pptx["FilePath"])
        print("Presentation opened successfully.")

    except Exception as e:
        print(f"An error occurred: {e}", file=sys.stderr)
        # log error to file for debugging
        with open("error_py.log", "a") as f:
            f.write(f"[ERROR] {datetime.datetime.now()}: An error occurred: {e}\n")
        # exit application with error code

    sys.exit(1)
